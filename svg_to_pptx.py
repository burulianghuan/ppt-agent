"""SVG 页面 -> PPTX。

策略：
- 优先用 cairosvg 把每页 SVG 渲成高清 PNG（Linux/Docker 环境，轻量）。
- 没有 cairosvg 时回退 headless 浏览器（Windows 本地用 Edge/Chrome）。
- PNG 作兜底位图，保证任何环境都能看。
- 把原始 SVG 作为矢量图注入到同一张图片上（svgBlip）。
  这样在 PowerPoint 2016+/365 里显示的是矢量，右键「转换为形状」即可编辑文字。

用法：
    python svg_to_pptx.py <svg目录> [输出.pptx]
    python svg_to_pptx.py outputs/7349401fc923
"""

from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# 16:9 幻灯片尺寸（EMU）：13.333in x 7.5in
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

# SVG 画布逻辑尺寸
SVG_W = 1280
SVG_H = 720
SCALE = 2  # 渲染倍率，PNG 兜底图清晰度

SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

# 浏览器候选（回退用）
_BROWSERS = (
    "msedge",
    "chromium",
    "chromium-browser",
    "google-chrome",
    "chrome",
)
_BROWSER_PATHS = (
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    "/usr/bin/chromium",
    "/usr/bin/chromium-browser",
    "/usr/bin/google-chrome",
)

_cairosvg_ok: bool | None = None


def cairosvg_available() -> bool:
    global _cairosvg_ok
    if _cairosvg_ok is None:
        try:
            import cairosvg  # noqa: F401

            _cairosvg_ok = True
        except Exception:
            _cairosvg_ok = False
    return _cairosvg_ok


def find_browser() -> str | None:
    for name in _BROWSERS:
        p = shutil.which(name)
        if p:
            return p
    for c in _BROWSER_PATHS:
        if Path(c).exists():
            return c
    return None


def _render_with_cairosvg(svg_path: Path, out_png: Path) -> None:
    import cairosvg

    cairosvg.svg2png(
        bytestring=svg_path.read_bytes(),
        write_to=str(out_png),
        output_width=SVG_W * SCALE,
        output_height=SVG_H * SCALE,
    )


def _render_with_browser(browser: str, svg_path: Path, out_png: Path) -> None:
    """用 headless 浏览器截图渲染 SVG -> PNG。"""
    svg = svg_path.read_text(encoding="utf-8")
    html = f"""<!doctype html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0}}html,body{{width:{SVG_W}px;height:{SVG_H}px;overflow:hidden}}
svg{{display:block;width:{SVG_W}px;height:{SVG_H}px}}</style></head>
<body>{svg}</body></html>"""

    with tempfile.TemporaryDirectory() as td:
        html_file = Path(td) / "page.html"
        html_file.write_text(html, encoding="utf-8")
        shot = Path(td) / "shot.png"
        cmd = [
            browser,
            "--headless=new",
            "--disable-gpu",
            "--no-sandbox",
            f"--window-size={SVG_W * SCALE},{SVG_H * SCALE}",
            f"--force-device-scale-factor={SCALE}",
            "--default-background-color=00000000",
            "--hide-scrollbars",
            f"--screenshot={shot}",
            html_file.as_uri(),
        ]
        subprocess.run(cmd, check=True, capture_output=True, timeout=120)
        if not shot.exists():
            raise RuntimeError(f"浏览器未生成截图: {svg_path.name}")
        out_png.write_bytes(shot.read_bytes())


def render_svg_to_png(svg_path: Path, out_png: Path) -> None:
    """渲染单页 SVG -> PNG，自动选择可用后端。"""
    if cairosvg_available():
        try:
            _render_with_cairosvg(svg_path, out_png)
            if out_png.exists() and out_png.stat().st_size > 0:
                return
        except Exception:
            pass  # cairosvg 失败则尝试浏览器

    browser = find_browser()
    if not browser:
        raise RuntimeError(
            "未找到 SVG 渲染器：请安装 cairosvg（Linux）或 Edge/Chrome（Windows）"
        )
    _render_with_browser(browser, svg_path, out_png)


def inject_svg(pptx_path: Path, mapping: list[tuple[int, Path]]) -> None:
    """把原始 SVG 作为矢量源注入每张图片（PowerPoint 显示矢量、可编辑）。

    mapping: [(slide_index, svg_path), ...]
    """
    from pptx.oxml.ns import qn

    prs = Presentation(str(pptx_path))
    for slide_idx, svg_path in mapping:
        slide = prs.slides[slide_idx]
        pic = None
        for shp in slide.shapes:
            if shp.shape_type == 13:  # PICTURE
                pic = shp
                break
        if pic is None:
            continue

        part = slide.part
        svg_bytes = svg_path.read_bytes()
        # 作为独立 part 加入并建立关系
        from pptx.opc.package import Part
        from pptx.opc.packuri import PackURI

        # 生成唯一 svg part 名
        n = slide_idx + 1
        svg_partname = PackURI(f"/ppt/media/vector{n}.svg")
        svg_part = Part(svg_partname, "image/svg+xml", package=part.package, blob=svg_bytes)
        rId = part.relate_to(
            svg_part,
            "http://schemas.microsoft.com/office/2007/relationships/hdphoto",
        )

        blipFill = pic._element.spPr.getparent().find(qn("p:blipFill"))
        if blipFill is None:
            blipFill = pic._element.find(qn("p:blipFill"))
        blip = pic._element.find(".//" + qn("a:blip"))
        if blip is None:
            continue
        # a:extLst / a:ext / asvg:svgBlip
        extLst = blip.makeelement(qn("a:extLst"), {})
        ext = blip.makeelement(qn("a:ext"), {"uri": "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"})
        svgBlip = blip.makeelement(
            "{%s}svgBlip" % SVG_NS,
            {qn("r:embed"): rId},
        )
        ext.append(svgBlip)
        extLst.append(ext)
        blip.append(extLst)

    prs.save(str(pptx_path))


def build_pptx(svg_dir: Path, out_pptx: Path) -> None:
    if not cairosvg_available() and not find_browser():
        raise RuntimeError(
            "未找到渲染器：请安装 cairosvg（Linux）或 Edge/Chrome（Windows）"
        )

    svgs = sorted(svg_dir.glob("*.svg"))
    if not svgs:
        raise RuntimeError(f"目录无 SVG: {svg_dir}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    mapping: list[tuple[int, Path]] = []
    tmpdir = Path(tempfile.mkdtemp(prefix="svg2pptx_"))

    for i, svg in enumerate(svgs):
        png = tmpdir / f"{i:02d}.png"
        print(f"[{i+1}/{len(svgs)}] 渲染 {svg.name} ...", flush=True)
        render_svg_to_png(svg, png)

        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=SLIDE_W, height=SLIDE_H)
        mapping.append((i, svg))

    prs.save(str(out_pptx))
    print("注入矢量 SVG（可编辑）...", flush=True)
    try:
        inject_svg(out_pptx, mapping)
        print("矢量注入完成。", flush=True)
    except Exception as e:
        print(f"矢量注入失败（PPTX 仍可用，为位图）: {e}", flush=True)

    shutil.rmtree(tmpdir, ignore_errors=True)
    print(f"完成: {out_pptx}  共 {len(svgs)} 页", flush=True)


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    src = Path(sys.argv[1]).resolve()
    out = Path(sys.argv[2]).resolve() if len(sys.argv) > 2 else src / "slides.pptx"
    build_pptx(src, out)
